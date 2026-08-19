# -*- coding: utf-8 -*-
"""Force PowerPoint to CLIP text to its text-frame (no paint-outside-box).

Sets on every text body:
  horzOverflow="clip"
  vertOverflow="clip"
  wrap="square"

Usage: python clip_pptx_text.py path/to/file.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE


def clip_shape_text(shape) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    tf = shape.text_frame
    try:
        tf.word_wrap = True
    except Exception:
        pass
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    try:
        body_pr = tf._txBody.bodyPr  # noqa: SLF001 — intentional OOXML touch
        body_pr.set("horzOverflow", "clip")
        body_pr.set("vertOverflow", "clip")
        body_pr.set("wrap", "square")
        return True
    except Exception:
        return False


def process(path: Path) -> int:
    prs = Presentation(str(path))
    n = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if clip_shape_text(shape):
                n += 1
            # tables / groups — best effort
            if shape.shape_type is not None and hasattr(shape, "shapes"):
                try:
                    for sub in shape.shapes:
                        if clip_shape_text(sub):
                            n += 1
                except Exception:
                    pass
    prs.save(str(path))
    return n


def main():
    if len(sys.argv) < 2:
        print("Usage: python clip_pptx_text.py file.pptx")
        sys.exit(2)
    path = Path(sys.argv[1])
    if not path.exists():
        print("Missing", path)
        sys.exit(1)
    n = process(path)
    print(f"Clipped text frames: {n} in {path.name}")


if __name__ == "__main__":
    main()
