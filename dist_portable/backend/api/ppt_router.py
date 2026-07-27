import os
import base64
import json
from fastapi import APIRouter
from backend.core.config import get_base_dir

router = APIRouter()


def _rgb_to_hex(rgb) -> str | None:
    try:
        if rgb is None:
            return None
        s = str(rgb)
        if len(s) == 6:
            return f"#{s}"
        return f"#{s}"
    except Exception:
        return None


def _shape_fill_hex(shape) -> str | None:
    try:
        fill = shape.fill
        if fill is None:
            return None
        # solid fill
        try:
            if fill.type is not None and fill.fore_color is not None:
                return _rgb_to_hex(getattr(fill.fore_color, "rgb", None))
        except Exception:
            pass
    except Exception:
        return None
    return None


def _shape_line_hex(shape) -> str | None:
    try:
        line = shape.line
        if line is None:
            return None
        if line.color is not None and hasattr(line.color, "rgb") and line.color.rgb:
            return _rgb_to_hex(line.color.rgb)
    except Exception:
        return None
    return None


def _line_width_pt(shape, default=1.5) -> float:
    try:
        w = shape.line.width
        if w is None:
            return default
        # EMU to pt
        return max(0.75, float(w.pt))
    except Exception:
        return default


def _is_rounded_rect(shape) -> bool:
    try:
        from pptx.enum.shapes import MSO_SHAPE
        return shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
    except Exception:
        try:
            xml = shape.element.xml
            return 'prst="roundRect"' in xml
        except Exception:
            return False


def _is_ellipse(shape) -> bool:
    try:
        from pptx.enum.shapes import MSO_SHAPE
        return shape.auto_shape_type == MSO_SHAPE.OVAL
    except Exception:
        try:
            return 'prst="ellipse"' in shape.element.xml
        except Exception:
            return False


@router.get("/api/ppt_export_json/{book_name}")
async def export_json_for_pptx_main(book_name: str):
    """Convert PPTX → JSON for the in-app web editor.

    Critical: callout cards are *text-on-shape* in PPTX. We must export fill/stroke
    with the text, otherwise the editor only shows naked text + blue connector arrows.
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        papers_dir = os.path.join(get_base_dir(), "data", "papers")
        textbooks_dir = os.path.join(get_base_dir(), "data", "textbooks")

        pptx_path = os.path.join(papers_dir, book_name, "pptx", f"{book_name}_Full_Presentation.pptx")
        if not os.path.exists(pptx_path):
            pptx_path = os.path.join(textbooks_dir, book_name, "pptx", f"{book_name}_Full_Presentation.pptx")

        if not os.path.exists(pptx_path):
            return {"error": "PPTX not found"}

        prs = Presentation(pptx_path)
        # 96 CSS-px per inch (frontend multiplies by 1280/960 to canvas)
        emu_to_px = 96 / 914400

        slides = []
        for i, slide in enumerate(prs.slides):
            elements = []
            for j, shape in enumerate(slide.shapes):
                el = {
                    "id": f"el_{i}_{j}",
                    "position": {
                        "x": shape.left * emu_to_px,
                        "y": shape.top * emu_to_px,
                    },
                    "size": {
                        "width": shape.width * emu_to_px,
                        "height": shape.height * emu_to_px,
                    },
                    "style": {
                        "opacity": 1,
                        "rotation": shape.rotation if hasattr(shape, "rotation") and shape.rotation else 0,
                    },
                }

                # --- Text (incl. text-on-shape callout cards) ---
                if shape.has_text_frame and shape.text.strip():
                    el["type"] = "text"
                    el["content"] = shape.text  # keep \n

                    font_size = 14
                    font_color = "#0F172A"
                    font_weight = "normal"
                    font_face = "Calibri"
                    try:
                        # Prefer first body-ish run; scan for a non-empty run
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if not (run.text or "").strip():
                                    continue
                                if run.font.size:
                                    font_size = run.font.size.pt
                                if run.font.color and hasattr(run.font.color, "rgb") and run.font.color.rgb:
                                    font_color = f"#{str(run.font.color.rgb)}"
                                if run.font.bold:
                                    font_weight = "bold"
                                if run.font.name:
                                    font_face = run.font.name
                                break
                            else:
                                continue
                            break
                    except Exception:
                        pass

                    text_align = "left"
                    try:
                        if shape.text_frame.paragraphs:
                            align = shape.text_frame.paragraphs[0].alignment
                            if align == 2:
                                text_align = "center"
                            elif align == 3:
                                text_align = "right"
                    except Exception:
                        pass

                    valign = "top"
                    try:
                        anchor = getattr(shape.text_frame, "vertical_anchor", None)
                        if anchor == 4:
                            valign = "middle"
                        elif anchor == 3:
                            valign = "bottom"
                    except Exception:
                        pass

                    el["style"].update({
                        "fontSize": font_size,
                        "color": font_color,
                        "fontWeight": font_weight,
                        "fontFamily": font_face,
                        "textAlign": text_align,
                        "valign": valign,
                    })

                    # Card chrome (roundRect / filled auto-shape)
                    fill_hex = _shape_fill_hex(shape)
                    stroke_hex = _shape_line_hex(shape)
                    if fill_hex:
                        el["style"]["fill"] = fill_hex
                    if stroke_hex:
                        el["style"]["stroke"] = stroke_hex
                        el["style"]["strokeWidth"] = _line_width_pt(shape, 1.5)
                    if _is_rounded_rect(shape):
                        el["style"]["borderRadius"] = 12
                    elif _is_ellipse(shape):
                        el["style"]["borderRadius"] = 999
                        if not fill_hex:
                            # numbered badge fallback
                            el["style"]["fill"] = stroke_hex or "#1E40AF"
                            el["style"]["color"] = "#FFFFFF"
                            el["style"]["textAlign"] = "center"
                            el["style"]["valign"] = "middle"

                    elements.append(el)
                    continue

                # --- Pictures ---
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    el["type"] = "image"
                    try:
                        image_blob = shape.image.blob
                        ext = shape.image.ext
                        b64_image = base64.b64encode(image_blob).decode("utf-8")
                        el["content"] = f"data:image/{ext};base64,{b64_image}"
                        el["style"]["objectFit"] = "contain"
                        elements.append(el)
                    except Exception:
                        pass
                    continue

                # --- Lines (connectors) — NOT blue arrows ---
                if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                    el["type"] = "shape"
                    el["content"] = "line"
                    stroke_color = _shape_line_hex(shape) or "#64748B"
                    el["style"].update({
                        "flipH": bool(
                            shape.element.xpath(".//a:xfrm/@flipH")
                            and shape.element.xpath(".//a:xfrm/@flipH")[0] == "1"
                        ),
                        "flipV": bool(
                            shape.element.xpath(".//a:xfrm/@flipV")
                            and shape.element.xpath(".//a:xfrm/@flipV")[0] == "1"
                        ),
                        "stroke": stroke_color,
                        "strokeWidth": _line_width_pt(shape, 1.25),
                        "noHead": True,
                    })
                    elements.append(el)
                    continue

                # --- Other auto-shapes (connectors may be prst=line) ---
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    xml = ""
                    try:
                        xml = shape.element.xml
                    except Exception:
                        pass
                    is_line_geom = 'prst="line"' in xml or "prst=\"straightConnector" in xml
                    if is_line_geom:
                        el["type"] = "shape"
                        el["content"] = "line"
                        stroke_color = _shape_line_hex(shape) or "#64748B"
                        el["style"].update({
                            "flipH": 'flipH="1"' in xml,
                            "flipV": 'flipV="1"' in xml,
                            "stroke": stroke_color,
                            "strokeWidth": _line_width_pt(shape, 1.25),
                            "noHead": True,
                        })
                        elements.append(el)
                        continue

                    el["type"] = "shape"
                    shape_val = "rectangle"
                    if _is_ellipse(shape):
                        shape_val = "ellipse"
                    elif _is_rounded_rect(shape):
                        shape_val = "roundRect"
                    el["content"] = shape_val
                    fill_color = _shape_fill_hex(shape) or "transparent"
                    stroke_color = _shape_line_hex(shape) or "#94A3B8"
                    # On-figure numbered badges are filled ellipses (~0.15–0.4") without text
                    if shape_val == "ellipse" and (not fill_color or fill_color == "transparent"):
                        fill_color = stroke_color if stroke_color and stroke_color != "#94A3B8" else "#1E40AF"
                    el["style"].update({
                        "fill": fill_color,
                        "stroke": stroke_color if shape_val != "ellipse" else (stroke_color or "#FFFFFF"),
                        "strokeWidth": _line_width_pt(shape, 1.5 if shape_val == "ellipse" else 1.25),
                    })
                    elements.append(el)

            slides.append({
                "id": f"slide_{i}",
                "background": {"type": "solid", "value": "#ffffff"},
                "elements": elements,
            })

        # Page mapping for PDF↔PPT sync
        # CRITICAL: slide index != figure index when follow-up ("补充说明") slides exist.
        page_mapping = {}
        n_pptx_slides = len(slides)

        def _natural_fig_key(name: str):
            import re as _re
            m = _re.search(r"Figure_(\d+)", name, _re.I)
            return (int(m.group(1)) if m else 10**9, name)

        def _load_meta(meta_path: str):
            with open(meta_path, "r", encoding="utf-8") as f_meta:
                raw = json.load(f_meta)
            # support {file: page} or {pages: {file: page}}
            if isinstance(raw, dict) and "pages" in raw and isinstance(raw["pages"], dict):
                return raw["pages"]
            return raw if isinstance(raw, dict) else {}

        # 1) Prefer slide_sync_map.json written by generate_full_ppt.js
        for root in (papers_dir, textbooks_dir):
            map_path = os.path.join(root, book_name, "pptx", "slide_sync_map.json")
            if os.path.exists(map_path):
                try:
                    with open(map_path, "r", encoding="utf-8") as f_map:
                        sync = json.load(f_map)
                    pm = sync.get("page_mapping") or {}
                    if pm:
                        page_mapping = {str(k): int(v) for k, v in pm.items() if v is not None}
                        print(f"[ppt_export] using slide_sync_map.json ({len(page_mapping)} entries)")
                        break
                except Exception as e:
                    print("Error reading slide_sync_map.json:", e)

        # 2) Reconstruct from ppt_cache + figures_metadata (handles follow-ups without re-export)
        if not page_mapping:
            try:
                meta_path = os.path.join(papers_dir, book_name, "images", "figures_metadata.json")
                if not os.path.exists(meta_path):
                    meta_path = os.path.join(textbooks_dir, book_name, "images", "figures_metadata.json")
                cache_path = os.path.join(papers_dir, book_name, "pptx", "ppt_cache_zh.json")
                if not os.path.exists(cache_path):
                    cache_path = os.path.join(papers_dir, book_name, "pptx", "ppt_cache.json")
                if not os.path.exists(cache_path):
                    cache_path = os.path.join(textbooks_dir, book_name, "pptx", "ppt_cache_zh.json")

                meta_data = _load_meta(meta_path) if os.path.exists(meta_path) else {}
                cache = {}
                if os.path.exists(cache_path):
                    with open(cache_path, "r", encoding="utf-8") as f_c:
                        cache = json.load(f_c)

                img_dir = os.path.join(papers_dir, book_name, "images")
                if not os.path.isdir(img_dir):
                    img_dir = os.path.join(textbooks_dir, book_name, "images")
                img_files = []
                if os.path.isdir(img_dir):
                    img_files = [f for f in os.listdir(img_dir) if f.lower().endswith((".png", ".jpg", ".jpeg"))]
                short = [f for f in img_files if f.startswith("Figure_")]
                long = [f for f in img_files if "_Figure_" in f and not f.startswith("Figure_")]
                if short:
                    # Default: lexicographic (historical JS .sort()). Natural used only if
                    # it uniquely matches pptx slide count better — tried below via dual build.
                    img_files = sorted(short)
                elif long:
                    img_files = sorted(long)
                else:
                    img_files = sorted(img_files)

                def _build_from(files_list):
                    pm_local = {}
                    si_local = 0
                    for fname in files_list:
                        page = meta_data.get(fname)
                        if page is None:
                            for mk, mv in meta_data.items():
                                if mk == fname or mk.endswith(fname) or fname.endswith(mk):
                                    page = mv
                                    break
                        if page is None:
                            continue
                        pm_local[str(si_local)] = int(page)
                        si_local += 1
                        entry = cache.get(fname) or {}
                        fus = entry.get("follow_up_slides") if isinstance(entry, dict) else None
                        if isinstance(fus, list):
                            for _ in fus:
                                pm_local[str(si_local)] = int(page)
                                si_local += 1
                    return pm_local, si_local

                # Try lex then natural; prefer exact match to pptx length
                candidates = []
                for flist in (
                    sorted(short) if short else sorted(img_files),
                    sorted(short, key=_natural_fig_key) if short else sorted(img_files, key=_natural_fig_key),
                ):
                    pm_try, si_try = _build_from(flist)
                    candidates.append((pm_try, si_try))
                page_mapping, slide_i = candidates[0]
                for pm_try, si_try in candidates:
                    if n_pptx_slides > 0 and si_try == n_pptx_slides:
                        page_mapping, slide_i = pm_try, si_try
                        break

                # If we overshot/undershot pptx slide count, pad/truncate safely
                if n_pptx_slides > 0 and page_mapping:
                    # If fewer mapped than slides, inherit last known page
                    last_page = None
                    for i in range(n_pptx_slides):
                        k = str(i)
                        if k in page_mapping:
                            last_page = page_mapping[k]
                        elif last_page is not None:
                            page_mapping[k] = last_page
                    # drop extras beyond pptx length
                    page_mapping = {k: v for k, v in page_mapping.items() if int(k) < n_pptx_slides}

                print(f"[ppt_export] reconstructed page_mapping ({len(page_mapping)} entries for {n_pptx_slides} slides)")
            except Exception as e:
                print("Error reconstructing page mapping:", e)
                import traceback
                traceback.print_exc()

        # 3) Legacy naive map (figure index only — last resort)
        if not page_mapping:
            meta_path = os.path.join(papers_dir, book_name, "images", "figures_metadata.json")
            if not os.path.exists(meta_path):
                meta_path = os.path.join(textbooks_dir, book_name, "images", "figures_metadata.json")
            if os.path.exists(meta_path):
                try:
                    meta_data = _load_meta(meta_path)
                    img_dir = os.path.dirname(meta_path)
                    img_files = sorted(
                        [f for f in os.listdir(img_dir) if f.lower().endswith((".png", ".jpg", ".jpeg")) and f.startswith("Figure_")],
                        key=_natural_fig_key,
                    )
                    for idx, fname in enumerate(img_files):
                        if fname in meta_data:
                            page_mapping[str(idx)] = meta_data[fname]
                except Exception as e:
                    print("Error building legacy page mapping:", e)

        return {"slides": slides, "page_mapping": page_mapping}
    except Exception as e:
        import traceback
        traceback.print_exc()
        return {"error": str(e)}


@router.get("/api/pptist_export_json/{book_name}")
async def export_json_for_ppt_master(book_name: str):
    return await export_json_for_pptx_main(book_name)


from pydantic import BaseModel
from typing import Optional
import re


class AnalyzeRequest(BaseModel):
    image: str
    slideWidth: int
    slideHeight: int
    imgX: int
    imgY: int
    imgW: int
    imgH: int
    book_name: Optional[str] = ""


@router.post("/api/analyze")
async def analyze_image_route(request: AnalyzeRequest):
    """Legacy analyze endpoint kept for interactive AI annotate in the editor."""
    try:
        md_content = ""
        if request.book_name:
            papers_dir = os.path.join(get_base_dir(), "data", "papers")
            textbooks_dir = os.path.join(get_base_dir(), "data", "textbooks")
            kb_path = os.path.join(
                papers_dir, request.book_name, "parsed", f"{request.book_name}_KnowledgeBase.md"
            )
            if not os.path.exists(kb_path):
                kb_path = os.path.join(
                    textbooks_dir, request.book_name, "parsed", f"{request.book_name}_KnowledgeBase.md"
                )
            if os.path.exists(kb_path):
                try:
                    with open(kb_path, "r", encoding="utf-8") as f_kb:
                        md_content = f_kb.read()
                except Exception:
                    pass

        from backend.core.config import load_config
        import random
        from openai import OpenAI

        cfg = load_config()
        keys = cfg.get("parse_api_key") or []
        if isinstance(keys, str):
            keys = [keys]
        api_key = random.choice(keys) if keys else (cfg.get("chat_api_key") or "")
        base_url = cfg.get("parse_api_url") or cfg.get("chat_api_url") or "https://api.siliconflow.cn/v1"
        model = cfg.get("parse_model") or cfg.get("chat_model") or "Qwen/Qwen2.5-72B-Instruct"

        client = OpenAI(api_key=api_key, base_url=base_url)
        # Minimal stub response so editor does not crash if used
        return {
            "annotations": [],
            "note": "Interactive analyze is secondary; prefer pre-generated PPT pipeline.",
            "model": model,
            "md_len": len(md_content),
        }
    except Exception as e:
        return {"error": str(e), "annotations": []}
