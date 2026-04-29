import os
import fitz # PyMuPDF
from fastapi import FastAPI, Request
from fastapi.responses import HTMLResponse, StreamingResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from pydantic import BaseModel
import io
import re
import json
import subprocess
import shutil
import asyncio
from openai import OpenAI
from fastapi import FastAPI, Request, UploadFile, File, BackgroundTasks, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import RedirectResponse

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
# 动态获取当前 main.py 所在的绝对路径，防止因为终端路径层级不对找不到 HTML 界面
BASE_APP_DIR = os.path.dirname(os.path.abspath(__file__))
templates = Jinja2Templates(directory=os.path.join(BASE_APP_DIR, "templates"))

BASE_APP_DIR = os.path.dirname(os.path.abspath(__file__))
templates = Jinja2Templates(directory=os.path.join(BASE_APP_DIR, "templates"))
app.mount("/static", StaticFiles(directory=os.path.join(BASE_APP_DIR, "static")), name="static")

import pythoncom
import win32com.client
from pptx import Presentation
import base64
import concurrent.futures

active_tasks = set()
# 使用多线程池实现并行处理能力，最大并发数4，搭配多秘钥轮询防止封禁
import random
task_queue_executor = concurrent.futures.ThreadPoolExecutor(max_workers=4)

def get_base_dir():
    # 获取此工程的外层根目录 (e:\workspace\ddl)
    return os.path.dirname(BASE_APP_DIR)

CONFIG_PATH = os.path.join(get_base_dir(), "config.json")

def load_config():
    if os.path.exists(CONFIG_PATH):
        with open(CONFIG_PATH, "r", encoding="utf-8") as f:
            cfg = json.load(f)
            # 向下兼容老版本的字符串秘钥
            if isinstance(cfg.get("parse_api_key"), str):
                cfg["parse_api_key"] = [k.strip() for k in cfg["parse_api_key"].split(",") if k.strip()]
            return cfg
    return {
        "parse_api_url": "https://api.siliconflow.cn/v1",
        "parse_api_key": [],
        "parse_model": "Qwen/Qwen3-VL-235B-A22B-Thinking",
        "chat_api_url": "https://api.siliconflow.cn/v1",
        "chat_api_key": "",
        "chat_model": "Qwen/Qwen3-VL-235B-A22B-Thinking"
    }

def save_config(new_config):
    with open(CONFIG_PATH, "w", encoding="utf-8") as f:
        json.dump(new_config, f, indent=4, ensure_ascii=False)

def scan_items(item_type="book"):
    items = []
    base_dir = get_base_dir()
    target_dir = "imports" if item_type == "book" else "papers"
    
    # We will just traverse the target directory
    store_dir = os.path.join(base_dir, target_dir)
    if os.path.exists(store_dir):
        for f in os.listdir(store_dir):
            if f.endswith(".pdf") and not f.endswith("_translated.pdf") and not f.endswith("_annotated.pdf"):
                b_name = os.path.splitext(f)[0]
                kb_file = os.path.join(store_dir, b_name, f"{b_name}_KnowledgeBase.md")
                pdf_file = os.path.join(store_dir, f)
                translated_pdf = os.path.join(store_dir, f"{b_name}_translated.pdf")
                annotated_pdf = os.path.join(store_dir, f"{b_name}_annotated.pdf")
                progress = ""
                # Check status
                if item_type == "book":
                    if os.path.exists(kb_file):
                        status = "ready"
                    else:
                        status = "processing" if f"books_{b_name}" in active_tasks else "interrupted"
                        
                        # Calculate progress
                        work_dir = os.path.join(store_dir, b_name)
                        temp_dir = os.path.join(work_dir, "temp_assets")
                        md_dir = os.path.join(work_dir, "markdown_parts")
                        
                        total_parts = 0
                        done_parts = 0
                        if os.path.exists(temp_dir):
                            total_parts = len([d for d in os.listdir(temp_dir) if os.path.isdir(os.path.join(temp_dir, d))])
                        if os.path.exists(md_dir):
                            done_parts = len([f for f in os.listdir(md_dir) if f.endswith(".md")])
                        
                        if total_parts > 0:
                            progress = f"{done_parts}/{total_parts}"
                        else:
                            progress = "抽取中"
                else:
                    pptx_path = os.path.join(store_dir, b_name, f"{b_name}_Full_Presentation.pptx")
                    status = "ready" if os.path.exists(pptx_path) else ("processing" if f"papers_{b_name}" in active_tasks else "interrupted")
                    
                items.append({
                    "name": b_name,
                    "pdf_path": pdf_file,
                    "translated_pdf_path": translated_pdf,
                    "annotated_pdf_path": annotated_pdf if os.path.exists(annotated_pdf) else "",
                    "kb_path": kb_file if os.path.exists(kb_file) else "",
                    "status": status,
                    "progress": progress,
                    "type": item_type
                })
    return items

def get_item_by_name(name):
    for b in scan_items("book") + scan_items("paper"):
        if b["name"] == name:
             return b
    return None

class ChatRequest(BaseModel):
    book_name: str
    message: str
    chat_history: list

@app.get("/api/config")
async def get_config():
    return load_config()

@app.post("/api/config")
async def update_config(req: Request):
    data = await req.json()
    save_config(data)
    return {"status": "success"}

@app.delete("/api/delete_target")
async def delete_target(name: str, type: str):
    base_dir = get_base_dir()
    target_dir = "imports" if type == "book" else "papers"
    store_dir = os.path.join(base_dir, target_dir)
    
    paths_to_delete = [
        os.path.join(store_dir, f"{name}.pdf"),
        os.path.join(store_dir, f"{name}_translated.pdf"),
        os.path.join(store_dir, f"{name}_annotated.pdf")
    ]
    
    import shutil
    for p in paths_to_delete:
        if os.path.exists(p):
            try: os.remove(p)
            except: pass
            
    work_dir = os.path.join(store_dir, name)
    if os.path.exists(work_dir):
        try: shutil.rmtree(work_dir, ignore_errors=True)
        except: pass
        
    return {"status": "success"}

active_tasks = set()

def run_builder_sync(pdf_path: str, book_name: str, item_type: str, prompt_type: str = "提示词汇总", ppt_mode: str = "creative"):
    task_id = f"{item_type}s_{book_name}"
    # active_tasks.add is now intentionally handled at submission site to ensure correct queued state
    try:
        import sys
        import shutil
        if item_type == "book":
            script_path = os.path.join(get_base_dir(), "universal_kb_builder.py")
            subprocess.run([sys.executable, script_path, pdf_path], check=True)
        else:
            target_dir = os.path.dirname(pdf_path) # e.g. E:\workspace\ddl\papers
            work_dir = os.path.join(target_dir, book_name) # e.g. E:\workspace\ddl\papers\FashionTex
            os.makedirs(work_dir, exist_ok=True)
            
            # Copy PDF into work_dir so pdf_annotator.py can find PDF and MD together
            work_pdf_path = os.path.join(work_dir, f"{book_name}.pdf")
            shutil.copy(pdf_path, work_pdf_path)
            
            # Step 1: Translate PDF
            print(f"\n========== Step 1: Translate PDF ==========")
            script_path = os.path.join(get_base_dir(), "tools", "paper_translator.py")
            out_translated = os.path.join(target_dir, f"{book_name}_translated.pdf")
            subprocess.run([sys.executable, script_path, pdf_path, out_translated], check=True)
            
            # Step 2: Extract Figures and Generate MD
            print(f"\n========== Step 2: Extract Figures & Gen Deep Parsing MD ==========")
            sys.path.insert(0, os.path.join(get_base_dir(), "standalone_pdf2ppt"))
            from project_manager import ProjectManager
            from llm_client import PaperReaderBot
            from prompts import get_stage1_prompt
            
            # 自动切分论文配图，为后续PPT使用提供素材
            pm = ProjectManager(base_dir=target_dir)
            pm.extract_semantic_figures(pdf_path, work_dir)
            
            cfg = load_config()
            
            # 动态密钥轮询逻辑以支持高并发吞吐
            parse_api_key_val = cfg.get("parse_api_key", [""])
            if isinstance(parse_api_key_val, str):
                parse_api_key_val = [k.strip() for k in parse_api_key_val.split(",") if k.strip()]
            
            paper_api_key_val = cfg.get("paper_api_key", "")
            valid_keys = []
            if isinstance(paper_api_key_val, str) and paper_api_key_val.strip():
                valid_keys = [k.strip() for k in paper_api_key_val.split(",") if k.strip()]
            elif isinstance(paper_api_key_val, list) and paper_api_key_val:
                valid_keys = [k for k in paper_api_key_val if k.strip()]
                
            if not valid_keys:
                valid_keys = [k for k in parse_api_key_val if k]
                
            api_key = random.choice(valid_keys) if valid_keys else ""
            
            base_url = cfg.get("paper_api_url", "https://api.siliconflow.cn/v1")
            model = cfg.get("paper_model", "Qwen/Qwen3-VL-235B-A22B-Thinking") 
            
              
            bot = PaperReaderBot(api_key=api_key, base_url=base_url, model_name=model)
            prompt = get_stage1_prompt(prompt_type)
            md_report = bot.get_stage1_md(pdf_path, prompt) 
            
            kb_file = os.path.join(work_dir, f"{book_name}_KnowledgeBase.md")
            with open(kb_file, "w", encoding="utf-8") as f:
                f.write(md_report)
            
            parsed_md = os.path.join(work_dir, f"输出结果_{book_name}.md")
            shutil.copy(kb_file, parsed_md)
            
            # Step 3: PPTX Compilation
            print(f"\\n========== Step 3: Compiling PPTX ==========")
            ppt_script = os.path.join(get_base_dir(), "standalone_pdf2ppt", "ppt_maker", "generate_full_ppt.js")
            out_ppt = os.path.join(work_dir, f"{book_name}_Full_Presentation.pptx")
            figures_dir = os.path.join(work_dir, "figures")
            
            import time
            for attempt in range(3):
                try:
                    subprocess.run(["node", ppt_script, kb_file, figures_dir, out_ppt, ppt_mode, api_key], cwd=os.path.join(get_base_dir(), "standalone_pdf2ppt", "ppt_maker"), check=True)
                    break
                except subprocess.CalledProcessError as e:
                    if attempt < 2:
                        print(f"PPT Compilation failed: {e}. Retrying ({attempt+2}/3)...")
                        time.sleep(2)
                    else:
                        raise e
            

            # Step 4: AI PDF Annotation
            print(f"\n========== Step 4: Generate Annotated PDF ==========")
            annotator_script = os.path.join(get_base_dir(), "tools", "pdf_annotator.py")
            subprocess.run([sys.executable, annotator_script, work_dir], check=True)
            
            ann_in_work = os.path.join(work_dir, f"{book_name}_annotated.pdf")
            ann_in_root = os.path.join(target_dir, f"{book_name}_annotated.pdf")
            if os.path.exists(ann_in_work):
                shutil.move(ann_in_work, ann_in_root)
            
            sys.path.pop(0)
    except subprocess.CalledProcessError as e:
        print(f"Error running processing: {e}")
    except Exception as e:
        import traceback
        traceback.print_exc()
    finally:
        active_tasks.discard(task_id)

@app.post("/api/upload")
async def upload_pdf(file: UploadFile = File(...)):
    return await handle_upload(file, "book")

@app.post("/api/upload_paper")
async def upload_paper(file: UploadFile = File(...), prompt_type: str = Form("提示词汇总"), ppt_mode: str = Form("creative")):
    return await handle_upload(file, "paper", prompt_type, ppt_mode)

async def handle_upload(file, item_type, prompt_type="提示词汇总", ppt_mode="creative"):
    target_dir = os.path.join(get_base_dir(), "imports" if item_type == "book" else "papers")
    os.makedirs(target_dir, exist_ok=True)
    
    filename = os.path.basename(file.filename) if file.filename else "unknown.pdf"
    
    # Strip spaces from filename to avoid Windows Path Error (e.g., "name  .pdf" -> "name.pdf")
    name_part, ext_part = os.path.splitext(filename)
    filename = name_part.strip() + ext_part.lower()

    if filename.endswith(".pdf.pdf"):
        filename = filename[:-4]
        
    pdf_path = os.path.join(target_dir, filename)
    content = await file.read()
    with open(pdf_path, "wb") as buffer:
        buffer.write(content)
        
    book_name = os.path.splitext(filename)[0]
    task_id = f"{item_type}s_{book_name}"
    if task_id not in active_tasks:
        active_tasks.add(task_id)
        if item_type == "paper":
            task_queue_executor.submit(run_builder_sync, pdf_path, book_name, item_type, prompt_type, ppt_mode)
        else:
            task_queue_executor.submit(run_builder_sync, pdf_path, book_name, item_type)
        
    return {"status": "processing", "book_name": book_name}

@app.post("/api/resume/{book_name}")
async def resume_task(book_name: str):
    imports_dir = os.path.join(get_base_dir(), "imports")
    pdf_path = os.path.join(imports_dir, f"{book_name}.pdf")
    if os.path.exists(pdf_path):
        task_id = f"books_{book_name}"
        if task_id not in active_tasks:
            active_tasks.add(task_id)
            task_queue_executor.submit(run_builder_sync, pdf_path, book_name, "book")
        return {"status": "processing"}
        
    papers_dir = os.path.join(get_base_dir(), "papers")
    pdf_path_paper = os.path.join(papers_dir, f"{book_name}.pdf")
    if os.path.exists(pdf_path_paper):
        task_id = f"papers_{book_name}"
        if task_id not in active_tasks:
            active_tasks.add(task_id)
            task_queue_executor.submit(run_builder_sync, pdf_path_paper, book_name, "paper", "提示词汇总", "creative")
        return {"status": "processing"}
        
    return {"status": "error", "message": "PDF not found"}

@app.get("/api/status/{item_type}/{book_name}")
async def check_status(item_type: str, book_name: str):
    target_dir = os.path.join(get_base_dir(), "imports" if item_type == "book" else "papers")
    if item_type == "book":
        kb_path = os.path.join(target_dir, book_name, f"{book_name}_KnowledgeBase.md")
        if os.path.exists(kb_path):
            return {"status": "ready"}
            
        status = "processing" if f"books_{book_name}" in active_tasks else "interrupted"
        work_dir = os.path.join(target_dir, book_name)
        temp_dir = os.path.join(work_dir, "temp_assets")
        md_dir = os.path.join(work_dir, "markdown_parts")
        
        total_parts = 0
        done_parts = 0
        if os.path.exists(temp_dir):
            total_parts = len([d for d in os.listdir(temp_dir) if os.path.isdir(os.path.join(temp_dir, d))])
        if os.path.exists(md_dir):
            done_parts = len([f for f in os.listdir(md_dir) if f.endswith(".md")])
            
        progress = f"{done_parts}/{total_parts}" if total_parts > 0 else "抽取中"
        return {"status": status, "progress": progress}
    else:
        pptx_path = os.path.join(target_dir, book_name, f"{book_name}_Full_Presentation.pptx")
        return {"status": "ready" if os.path.exists(pptx_path) else "processing"}

@app.get("/", response_class=HTMLResponse)
async def index(request: Request):
    books = scan_items("book")
    papers = scan_items("paper")
    return templates.TemplateResponse("index.html", {"request": request, "books": books, "papers": papers})

@app.get("/chat/{book_name}", response_class=HTMLResponse)
async def chat_page(request: Request, book_name: str):
    book = get_item_by_name(book_name)
    if not book:
        return "Book not found", 404
    return templates.TemplateResponse("chat.html", {"request": request, "book_name": book_name, "is_paper": book["type"] == "paper"})

@app.get("/cover/{book_name}")
async def get_cover(book_name: str):
    book = get_item_by_name(book_name)
    if not book: return "404", 404
    
    # 借助用 PyMuPDF 瞬间截取 PDF 第一页作为首页的封面图片展示
    doc = fitz.open(book["pdf_path"])
    page = doc.load_page(0)
    pix = page.get_pixmap(matrix=fitz.Matrix(0.5, 0.5))
    img_data = pix.tobytes("png")
    doc.close()
    
    return StreamingResponse(io.BytesIO(img_data), media_type="image/png")

@app.get("/pdf/{book_name}")
async def get_pdf(book_name: str):
    item = get_item_by_name(book_name)
    if not item: return "404", 404
    return FileResponse(item["pdf_path"], media_type="application/pdf")

@app.get("/pdf_translated/{book_name}")
async def get_pdf_translated(book_name: str):
    item = get_item_by_name(book_name)
    if not item or not os.path.exists(item.get("translated_pdf_path", "")) : return "404", 404
    return FileResponse(item["translated_pdf_path"], media_type="application/pdf")

@app.get("/pdf_annotated/{book_name}")
async def get_pdf_annotated(book_name: str):
    item = get_item_by_name(book_name)
    if not item or not os.path.exists(item.get("annotated_pdf_path", "")) : return "404", 404
    return FileResponse(item["annotated_pdf_path"], media_type="application/pdf")

# ================= PPT 编辑核心框架 =================
@app.get("/ppt_editor/{book_name}")
async def ppt_editor_page(request: Request, book_name: str):
    # 重定向到 ppt-master Vue 项目 (Vite 运行在 8081 端口，由于处于 WSL 环境，使用 localhost 对 Windows 端口转发兼容性最好)
    return RedirectResponse(f"http://localhost:8081/?book={book_name}")

@app.get("/api/ppt_export_json/{book_name}")
async def export_json_for_pptx_main(book_name: str):
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        import base64
        import os
        
        # Check both the specific paper imports and book imports
        imports_dir = os.path.join(get_base_dir(), "imports", book_name)
        papers_dir = os.path.join(get_base_dir(), "papers", book_name)
        
        pptx_path = os.path.join(papers_dir, f"{book_name}_Full_Presentation.pptx")
        if not os.path.exists(pptx_path):
            pptx_path = os.path.join(imports_dir, f"{book_name}_Full_Presentation.pptx")
            
        if not os.path.exists(pptx_path):
            return {"error": "PPTX not found"}
            
        prs = Presentation(pptx_path)
        emu_to_px = 96 / 914400
        
        slides = []
        for i, slide in enumerate(prs.slides):
            elements = []
            for j, shape in enumerate(slide.shapes):
                el = {
                    "id": f"el_{i}_{j}",
                    "position": {"x": shape.left * emu_to_px, "y": shape.top * emu_to_px},
                    "size": {"width": shape.width * emu_to_px, "height": shape.height * emu_to_px},
                    "style": {
                        "opacity": 1,
                        "rotation": shape.rotation if hasattr(shape, 'rotation') and shape.rotation else 0
                    }
                }
                
                if shape.has_text_frame and shape.text.strip():
                    el["type"] = "text"
                    el["content"] = shape.text
                    
                    font_size = 18
                    font_color = "#000000"
                    font_weight = "normal"
                    
                    try:
                        # Extract exact font style from the first run
                        if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                            run = shape.text_frame.paragraphs[0].runs[0]
                            if run.font.size:
                                font_size = run.font.size.pt
                            if run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                font_color = f"#{str(run.font.color.rgb)}"
                            if run.font.bold:
                                font_weight = "bold"
                    except:
                        pass
                        
                    
                    text_align = "left"
                    try:
                        if shape.text_frame.paragraphs:
                            align = shape.text_frame.paragraphs[0].alignment
                            if align == 2:
                                text_align = "center"
                            elif align == 3:
                                text_align = "right"
                    except:
                        pass
                    
                    valign = "top"
                    try:
                        anchor = getattr(shape.text_frame, 'vertical_anchor', None)
                        if anchor == 4:
                            valign = "middle"
                        elif anchor == 3:
                            valign = "bottom"
                    except:
                        pass
                    
                    el["style"].update({
                        "fontSize": font_size,
                        "color": font_color,
                        "fontWeight": font_weight,
                        "textAlign": text_align,
                        "valign": valign
                    })
                    elements.append(el)
                    
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    el["type"] = "image"
                    try:
                        image_blob = shape.image.blob
                        ext = shape.image.ext
                        b64_image = base64.b64encode(image_blob).decode('utf-8')
                        el["content"] = f"data:image/{ext};base64,{b64_image}"
                        el["style"]["objectFit"] = "contain"
                        elements.append(el)
                    except:
                        pass
                    
                elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
                    el["type"] = "shape"
                    el["content"] = "arrow" # Map lines to arrows as requested
                    stroke_color = "#3b82f6"
                    try:
                        if shape.line.color and hasattr(shape.line.color, 'rgb') and shape.line.color.rgb:
                            stroke_color = f"#{str(shape.line.color.rgb)}"
                    except:
                        pass
                    el["style"].update({
                        "flipH": bool(shape.element.xpath('.//a:xfrm/@flipH') and shape.element.xpath('.//a:xfrm/@flipH')[0] == '1'),
                        "flipV": bool(shape.element.xpath('.//a:xfrm/@flipV') and shape.element.xpath('.//a:xfrm/@flipV')[0] == '1'),
                        "stroke": stroke_color,
                        "strokeWidth": 2
                    })
                    elements.append(el)
                    
                elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    el["type"] = "shape"
                    shape_val = "rectangle"
                    try:
                        ast = getattr(shape, 'auto_shape_type', None)
                        if ast in (33, 34, 35, 36) or ast == 9: # Arrow types or Line
                            shape_val = "arrow"
                    except:
                        if hasattr(shape, 'element') and 'prst="line"' in shape.element.xml:
                            shape_val = "line"
                    
                    if shape_val in ("line", "arrow"):
                        el["style"]["flipH"] = bool(shape.element.xpath('.//a:xfrm/@flipH') and shape.element.xpath('.//a:xfrm/@flipH')[0] == '1')
                        el["style"]["flipV"] = bool(shape.element.xpath('.//a:xfrm/@flipV') and shape.element.xpath('.//a:xfrm/@flipV')[0] == '1')
                    el["content"] = shape_val
                    fill_color = "transparent"
                    stroke_color = "#3b82f6"
                    try:
                        if shape.fill and shape.fill.solid() and shape.fill.fore_color and hasattr(shape.fill.fore_color, 'rgb') and shape.fill.fore_color.rgb:
                            fill_color = f"#{str(shape.fill.fore_color.rgb)}"
                    except:
                        pass
                    el["style"].update({
                        "fill": fill_color,
                        "stroke": stroke_color,
                        "strokeWidth": 2
                    })
                    elements.append(el)
                    
            slides.append({
                "id": f"slide_{i}",
                "background": {"type": "solid", "value": "#ffffff"},
                "elements": elements
            })
            
        return {"slides": slides}
    except Exception as e:
        import traceback
        traceback.print_exc()
        return {"error": str(e)}

@app.get("/api/pptist_export_json/{book_name}")
async def export_json_for_ppt_master(book_name: str):
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        import base64
        import os
        
        imports_dir = os.path.join(get_base_dir(), "imports", book_name)
        papers_dir = os.path.join(get_base_dir(), "papers", book_name)
        
        pptx_path = os.path.join(papers_dir, f"{book_name}_Full_Presentation.pptx")
        if not os.path.exists(pptx_path):
            pptx_path = os.path.join(imports_dir, f"{book_name}_Full_Presentation.pptx")
            
        if not os.path.exists(pptx_path):
            return {"error": "PPTX not found"}
            
        prs = Presentation(pptx_path)
        # Convert EMU to px for 1000 width viewport (PPTist standard)
        emu_to_px = 1000 / 9144000
        
        slides = []
        for i, slide in enumerate(prs.slides):
            elements = []
            for j, shape in enumerate(slide.shapes):
                base_el = {
                    "id": f"el_{i}_{j}",
                    "left": shape.left * emu_to_px,
                    "top": shape.top * emu_to_px,
                    "width": shape.width * emu_to_px,
                    "height": shape.height * emu_to_px,
                    "rotate": shape.rotation if hasattr(shape, 'rotation') and shape.rotation else 0
                }
                
                fill_color = "transparent"
                try:
                    if hasattr(shape, 'fill') and shape.fill.type == 1:
                        if hasattr(shape.fill.fore_color, 'rgb') and shape.fill.fore_color.rgb:
                            fill_color = f"#{str(shape.fill.fore_color.rgb)}"
                        else:
                            fill_color = "#3b82f6"
                except:
                    pass

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_blob = shape.image.blob
                        ext = shape.image.ext
                        b64_image = base64.b64encode(image_blob).decode('utf-8')
                        el = {
                            **base_el,
                            "type": "image",
                            "src": f"data:image/{ext};base64,{b64_image}",
                            "fixedRatio": False
                        }
                        elements.append(el)
                    except:
                        pass
                        
                elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
                    stroke_color = "#3b82f6"
                    try:
                        if hasattr(shape, 'line') and shape.line.color and hasattr(shape.line.color, 'rgb') and shape.line.color.rgb:
                            stroke_color = f"#{str(shape.line.color.rgb)}"
                    except:
                        pass
                        
                    el = {
                        "id": base_el["id"],
                        "left": base_el["left"],
                        "top": base_el["top"],
                        "width": 3,
                        "type": "line",
                        "start": [0, 0],
                        "end": [base_el["width"], base_el["height"]],
                        "style": "solid",
                        "color": stroke_color,
                        "points": ["", "arrow"]
                    }
                    elements.append(el)

                elif shape.has_text_frame and shape.text.strip():
                    font_size = 18
                    font_color = "#000000"
                    is_bold = False
                    
                    try:
                        if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                            run = shape.text_frame.paragraphs[0].runs[0]
                            if run.font.size:
                                font_size = run.font.size.pt
                            if run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                font_color = f"#{str(run.font.color.rgb)}"
                            if run.font.bold:
                                is_bold = True
                    except:
                        pass
                        
                    font_size_px = int(font_size * 1.3888) # Scale to 1000-Viewport
                    text_html = shape.text.replace('\n', '<br>')
                    html_content = f"<p style='color: {font_color}; font-size: {font_size_px}px; font-weight: {'bold' if is_bold else 'normal'}; margin: 0; padding: 0;'>{text_html}</p>"
                    
                    el = {
                        **base_el,
                        "left": base_el["left"] - 10,
                        "top": base_el["top"] - 10,
                        "width": base_el["width"] + 20,
                        "height": base_el["height"] + 20,
                        "type": "text",
                        "content": html_content,
                        "defaultFontName": "Arial",
                        "defaultColor": font_color,
                        "lineHeight": 1.2,
                        "fill": fill_color
                    }
                    elements.append(el)
                    
                else:
                    # Check if it's actually a line masquerading as a generic shape
                    xml_str = shape.element.xml
                    is_custom_line = 'prst="line"' in xml_str
                    
                    if is_custom_line:
                        # Extract flipV and flipH
                        flip_v = 'flipV="1"' in xml_str
                        flip_h = 'flipH="1"' in xml_str
                        
                        start_pt = [0, 0]
                        end_pt = [base_el["width"], base_el["height"]]
                        if flip_v and not flip_h:
                            start_pt = [0, base_el["height"]]
                            end_pt = [base_el["width"], 0]
                        elif flip_h and not flip_v:
                            start_pt = [base_el["width"], 0]
                            end_pt = [0, base_el["height"]]
                        elif flip_h and flip_v:
                            start_pt = [base_el["width"], base_el["height"]]
                            end_pt = [0, 0]
                            
                        # To keep LINE from getting height/rotate, we filter base_el
                        line_el = { k: v for k, v in base_el.items() if k not in ["height", "rotate"] }
                        el = {
                            **line_el,
                            "type": "line",
                            "start": start_pt,
                            "end": end_pt,
                            "width": 3,
                            "style": "solid",
                            "color": fill_color if fill_color != "transparent" else "#3b82f6",
                            "points": ["", ""]
                        }
                        elements.append(el)
                    else:
                        ast = 1 # RECTANGLE
                        try:
                            ast = getattr(shape, 'auto_shape_type', 1)
                            # Handle enum value if it is not integer directly
                            if hasattr(ast, 'value'):
                                ast = ast.value
                        except:
                            pass
                            
                        view_box = [200, 200]
                        path = 'M 0 0 L 200 0 L 200 200 L 0 200 Z' # Rect by default
                        
                        if ast == 7: # ISOSCELES_TRIANGLE
                            path = 'M 100 0 L 0 200 L 200 200 Z'
                            
                        el = {
                            **base_el,
                            "type": "shape",
                            "viewBox": view_box,
                            "path": path,
                            "fill": fill_color if fill_color != "transparent" else "#3b82f6",
                            "fixedRatio": False
                        }
                        elements.append(el)
                    
            slides.append({
                "id": f"slide_{i}",
                "elements": elements,
                "background": {
                    "type": "solid",
                    "color": "#ffffff"
                }
            })
            
        return {"slides": slides}
    except Exception as e:
        import traceback
        traceback.print_exc()
        return {"error": str(e)}


# 核心：参考 Cherry Studio 的前置预处理器 (Preprocess => Retrieve)
def simple_rag_search(kb_path, query):
    '''
    这里借鉴 Cherry Studio 结构，将巨无霸 markdown 拆块。
    Cherry 原生通过 embedjs 转化为向量；
    这里为了不增加复杂的 chroma / faiss 本体环境依赖，我们采用敏捷“轻量级文本 TF-IDF / 关键词交集匹配法”替代。
    由于 universal_kb_builder 的生成物质量极高，即便是轻量级匹配也能达到恐怖的精度。
    '''
    with open(kb_path, "r", encoding="utf-8") as f:
         content = f.read()
         
    # 根据页码标志强行切碎！
    chunks = re.split(r'\n##\s+|\n###\s+', content)
    
    # 极其粗暴却有效的本地化分词打分器
    stopwords = ["的", "了", "吗", "呢", "什么是", "怎么", "?", "？"]
    keywords = set([w for w in query if w not in stopwords])
    if not keywords: keywords = set(query)
    
    scored_chunks = []
    for chunk in chunks:
        score = sum([2 for kw in keywords if kw in chunk])
        if score > 0:
            scored_chunks.append((score, chunk))
            
    scored_chunks.sort(key=lambda x: x[0], reverse=True)
    
    # 将相关度最靠前的内容缝在一起，限制喂给大模型的 Token，预防超流
    top_chunks = [c[1][:1000] for c in scored_chunks[:3]]
    return "\n======\n".join(top_chunks)

@app.post("/api/chat")
async def chat_api(req: ChatRequest):
    book = get_item_by_name(req.book_name)
    if not book: return {"reply": "Book not found"}
    
    try:
        # Step 1: 文本检索召回 (Retrieve)
        # 模仿 Cherry 的 \`embedjs.search()\` 行为
        context = simple_rag_search(book["kb_path"], req.message)
        
        # Step 2: 组装融合知识 (Augmented)
        sys_prompt = f"""你是【{req.book_name}】专属私教助教。
请严格基于底下提取的【知识库内部切片信息】来回答用户，如果里面没提到，请不要生搬硬造。

【提取到的教材切片信息】：
{context}
"""
        
        messages = [{"role": "system", "content": sys_prompt}]
        for hist in req.chat_history:
             # 把以前的历史塞进去保证多轮聊天的连贯
            messages.append({"role": hist["role"], "content": hist["content"]})
        messages.append({"role": "user", "content": req.message})
        
        # 实时加载设置界面刚配好的引擎和秘钥
        cfg = load_config()
        chat_client = OpenAI(api_key=cfg["chat_api_key"], base_url=cfg["chat_api_url"])
        
        # Step 3: 原生语言生成器 (Generation)
        response = chat_client.chat.completions.create(
            model=cfg.get("chat_model", "Qwen/Qwen3-VL-235B-A22B-Thinking"),
            messages=messages,
            max_tokens=2048,
            temperature=0.7
        )
        return {"reply": response.choices[0].message.content}
    except Exception as e:
        return {"reply": f"API调用遇到错误：{str(e)}"}

# ==================== 提示词管理 API ====================
class PromptSaveRequest(BaseModel):
    content: str

@app.get("/api/prompts")
async def list_prompts():
    prompt_dir = os.path.join(get_base_dir(), "standalone_pdf2ppt", "prompts")
    if not os.path.exists(prompt_dir):
        os.makedirs(prompt_dir)
    files = [f for f in os.listdir(prompt_dir) if f.endswith('.md')]
    names = [os.path.splitext(f)[0] for f in files]
    return {"status": "success", "prompts": names}

@app.get("/api/prompts/{prompt_name}")
async def get_prompt(prompt_name: str):
    prompt_dir = os.path.join(get_base_dir(), "standalone_pdf2ppt", "prompts")
    path = os.path.join(prompt_dir, f"{prompt_name}.md")
    if not os.path.exists(path):
        return {"status": "error", "message": "Prompt not found"}
    with open(path, "r", encoding="utf-8") as f:
        return {"status": "success", "content": f.read()}

@app.post("/api/prompts/{prompt_name}")
async def save_prompt(prompt_name: str, req: PromptSaveRequest):
    prompt_name = os.path.basename(prompt_name)
    prompt_dir = os.path.join(get_base_dir(), "standalone_pdf2ppt", "prompts")
    if not os.path.exists(prompt_dir):
        os.makedirs(prompt_dir)
    path = os.path.join(prompt_dir, f"{prompt_name}.md")
    with open(path, "w", encoding="utf-8") as f:
        f.write(req.content)
    return {"status": "success"}

@app.delete("/api/prompts/{prompt_name}")
async def delete_prompt(prompt_name: str):
    prompt_name = os.path.basename(prompt_name)
    prompt_dir = os.path.join(get_base_dir(), "standalone_pdf2ppt", "prompts")
    path = os.path.join(prompt_dir, f"{prompt_name}.md")
    if os.path.exists(path):
        os.remove(path)
        return {"status": "success"}
    return {"status": "error", "message": "Prompt not found"}

if __name__ == "__main__":
    import uvicorn
    import webbrowser
    import threading
    import time
    
    # 自动弹窗机制：等服务器启动1.5秒后，自动唤起浏览器
    def open_browser():
        time.sleep(1.5)
        webbrowser.open("http://127.0.0.1:8899/")
        
    threading.Thread(target=open_browser, daemon=True).start()
    
    # 我们换一个冷门的 8899 端口
    uvicorn.run(app, host="127.0.0.1", port=8899)
