# -*- coding: utf-8 -*-
import json
from pathlib import Path
from pptx import Presentation
from PIL import Image

p = list(Path("data/papers").glob("Y-zipper*"))[0]
print("dir:", p.name)

pptx = [x for x in (p / "pptx").glob("*_Full_Presentation.pptx") if "EN" not in x.name and "__tmp" not in x.name][0]
prs = Presentation(str(pptx))
print("n_slides:", len(prs.slides), "size_mb:", round(pptx.stat().st_size / 1e6, 2))

for i, slide in enumerate(prs.slides):
    if i >= 8:
        break
    texts = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text.strip():
            t = sh.text.strip().replace("\n", " | ")
            if len(t) > 140:
                t = t[:140] + "..."
            texts.append(t)
    print(f"--- slide {i+1} ---")
    for t in texts[:8]:
        print(" ", t)

cache = json.load(open(p / "pptx" / "ppt_cache_zh.json", encoding="utf-8"))
print("\ncache figures:", len(cache))
fus = sum(len(v.get("follow_up_slides") or []) for v in cache.values() if isinstance(v, dict))
print("followups total:", fus)

# natural order print
import re
def nk(k):
    m = re.search(r"Figure_(\d+)", k)
    return int(m.group(1)) if m else 999

for k in sorted(cache.keys(), key=nk):
    v = cache[k]
    if not isinstance(v, dict):
        continue
    title = (v.get("slide_title") or "")[:50]
    anns = v.get("annotations") or []
    labels = [((a.get("label") or "")[:20]) for a in anns[:4]]
    print(f"{k}: ann={len(anns)} fu={len(v.get('follow_up_slides') or [])} title={title!r} labels={labels}")

# image quality heuristic: very tall crops may include body text
print("\n--- image aspect ---")
img = p / "images"
for f in sorted(img.glob("Figure_*.png"), key=lambda x: nk(x.name)):
    im = Image.open(f)
    ar = im.size[1] / max(im.size[0], 1)
    flag = " TALL?" if ar > 1.15 else ""
    print(f"{f.name}: {im.size} ar={ar:.2f}{flag}")
