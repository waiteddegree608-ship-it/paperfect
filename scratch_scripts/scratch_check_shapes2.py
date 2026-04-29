import glob
from pptx import Presentation

p = glob.glob('*_Full_Presentation.pptx')[0]
prs = Presentation(p)

for i, sl in enumerate(list(prs.slides)[:3]):
    print(f"Slide {i}:")
    for s in sl.shapes:
        has_line = hasattr(s, 'element') and 'prst="line"' in s.element.xml
        has_text = s.has_text_frame
        text = s.text.strip() if has_text else ''
        print(s.shape_type, "has_line:", has_line, "has_text:", has_text, "text:", repr(text))
