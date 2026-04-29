import glob
from pptx import Presentation

p = glob.glob('*_Full_Presentation.pptx')[0]
prs = Presentation(p)

for i, sl in enumerate(prs.slides[:3]):
    print(f"Slide {i}:")
    for s in sl.shapes:
        has_line = hasattr(s, 'element') and 'prst="line"' in s.element.xml
        has_triangle = hasattr(s, 'element') and 'prst="triangle"' in s.element.xml
        text = s.text if s.has_text_frame else ''
        print(s.shape_type, "has_line:", has_line, "has_triangle:", has_triangle, "text:", repr(text[:20]))
